#!/usr/bin/env python3
"""
Generate Phonetic Contrastive Model presentation (PPTX).

Produces: ard/Phonetic-Contrastive-Model-Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn, nsmap
from lxml import etree
from pathlib import Path

OUTPUT = Path(__file__).resolve().parent / "Phonetic-Contrastive-Model-Presentation.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Colours ──────────────────────────────────────────────────────────────────
DARK_BG      = RGBColor(0x0F, 0x17, 0x2A)
CARD_BG      = RGBColor(0x1A, 0x25, 0x3C)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xBB, 0xBB, 0xBB)
ACCENT_BLUE  = RGBColor(0x4E, 0xC5, 0xF1)
ACCENT_GREEN = RGBColor(0x4E, 0xC9, 0xA0)
ACCENT_ORANGE= RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED   = RGBColor(0xE7, 0x4C, 0x3C)
MUTED_BLUE   = RGBColor(0x34, 0x98, 0xDB)
CODE_BG      = RGBColor(0x12, 0x1A, 0x2D)

blank_layout = prs.slide_layouts[6]


def set_slide_bg(slide, color=DARK_BG):
    """Set slide background via the proper <p:bg> element — universal rendering."""
    cSld = slide._element.find(qn("p:cSld"))
    # Remove any existing <p:bg>
    for old in cSld.findall(qn("p:bg")):
        cSld.remove(old)
    # Build <p:bg><p:bgPr><a:solidFill><a:srgbClr val="..."/></a:solidFill>
    #           <a:effectLst/></p:bgPr></p:bg>
    bg = etree.SubElement(cSld, qn("p:bg"))
    cSld.remove(bg)
    cSld.insert(0, bg)  # <p:bg> must be first child of <p:cSld>
    bgPr = etree.SubElement(bg, qn("p:bgPr"))
    solidFill = etree.SubElement(bgPr, qn("a:solidFill"))
    etree.SubElement(solidFill, qn("a:srgbClr"), val=str(color))
    etree.SubElement(bgPr, qn("a:effectLst"))


def add_text_box(slide, left, top, width, height, text,
                 font_size=14, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    if font_name:
        p.font.name = font_name
    return tf


def add_bullet_frame(slide, left, top, width, height, items,
                     font_size=16, color=WHITE, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return tf


def add_code_box(slide, left, top, width, height, text, font_size=11):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    # dark code background via the shape's fill
    fill = txBox.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x12, 0x1A, 0x2D)

    tf = txBox.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(0xA0, 0xD9, 0x68)
        p.font.name = "Menlo"
        p.space_after = Pt(1)
    return tf


def add_slide_number(slide, num, total):
    add_text_box(slide, Inches(12.4), Inches(7.0), Inches(0.8), Inches(0.4),
                 f"{num}/{total}", font_size=10, color=LIGHT_GRAY,
                 alignment=PP_ALIGN.RIGHT)


def add_header(slide, title, subtitle="", slide_num=0, total=0):
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.7),
                 title, font_size=32, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, Inches(0.6), Inches(0.95), Inches(10), Inches(0.4),
                     subtitle, font_size=14, color=LIGHT_GRAY)
    if slide_num:
        add_slide_number(slide, slide_num, total)


TOTAL_SLIDES = 15

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(1.0),
             "Phonetic Contrastive Model", font_size=44, bold=True, color=WHITE)
add_text_box(slide, Inches(0.8), Inches(2.9), Inches(11), Inches(0.8),
             "Character-Level Siamese Bi-Encoder for Roman Urdu Spelling Normalisation",
             font_size=22, color=ACCENT_BLUE)
add_text_box(slide, Inches(0.8), Inches(4.2), Inches(11), Inches(0.5),
             "A learned replacement for rule-based spelling correction in the ASR pipeline",
             font_size=16, color=LIGHT_GRAY)
add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.5),
             "BitLogix  ·  July 2026", font_size=14, color=LIGHT_GRAY)
add_slide_number(slide, 1, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_header(slide, "The Problem: Roman Urdu Has No Standard Spelling",
           "Same word, dozens of valid spellings — a normalisation nightmare", 2, TOTAL_SLIDES)

add_bullet_frame(slide, Inches(0.6), Inches(1.7), Inches(5.8), Inches(2.5), [
    "▸ Roman Urdu is written in Latin script with no official orthography",
    "▸ ASR output (Hindi → Roman) produces phonetically correct but inconsistent spellings",
    "▸ Downstream systems expect canonical forms",
    "▸ Previous rule-based resolver: ~30% recall, O(N²) maintenance",
], font_size=16, color=WHITE)

# Example table
add_text_box(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.4),
             "Real-World Spelling Variants", font_size=18, bold=True, color=ACCENT_BLUE)

examples = [
    ("Canonical", "Observed Variants"),
    ("Siddiqui", "siddiqee, siddiqi, siddique, sidiqui"),
    ("Chughtai", "chugataai, chugatai, chughtaai"),
    ("operation", "apareshan, opreshan, aperation"),
    ("hospital", "haspatal, hispatal, hospitaal"),
]
y = Inches(2.1)
for canon, variants in examples:
    is_header = canon == "Canonical"
    c = ACCENT_BLUE if is_header else ACCENT_GREEN
    add_text_box(slide, Inches(7.0), y, Inches(2.0), Inches(0.35),
                 canon, font_size=14, bold=is_header, color=c)
    add_text_box(slide, Inches(9.0), y, Inches(4.0), Inches(0.35),
                 variants, font_size=13, color=WHITE if is_header else LIGHT_GRAY)
    y += Inches(0.4)

add_text_box(slide, Inches(0.6), Inches(5.0), Inches(12), Inches(0.8),
             "Goal: A model that generalises to UNSEEN spellings with built-in safety against false corrections",
             font_size=18, bold=True, color=ACCENT_ORANGE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Solution Overview
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_header(slide, "Solution: Metric Learning on Characters",
           "Learn an embedding space where phonetically equivalent words cluster together", 3, TOTAL_SLIDES)

add_bullet_frame(slide, Inches(0.6), Inches(1.7), Inches(5.8), Inches(4.5), [
    "▸ Siamese bi-encoder: same network embeds variants AND canonicals",
    "▸ Character-level: signal is phonetic (z↔j, aa↔a, ee↔i), not semantic",
    "▸ Contrastive training: InfoNCE loss pulls variants toward their canonical",
    "▸ Pre-computed canonical index: inference = 1 forward pass + 1 matmul",
    "▸ Abstain threshold (0.90): uncertain → leave word unchanged",
    "▸ Result: ~97% held-out recall vs resolver's ~30%",
], font_size=16, color=WHITE)

add_code_box(slide, Inches(7.0), Inches(1.7), Inches(5.8), Inches(4.0),
             '# Inference — it\'s this simple\n'
             'from phonetic_contrastive_model.corrector import \\\n'
             '    PhoneticContrastiveCorrector\n\n'
             'c = PhoneticContrastiveCorrector.load()\n\n'
             'c.resolve_word("chugataai")   # → "Chughtai"\n'
             'c.resolve_word("siddiqee")    # → "Siddiqui"\n'
             'c.resolve_word("apareshan")   # → "apareshan"  (abstained)\n'
             'c.resolve_word("area")        # → "area"  (already known)\n\n'
             '# Add new canonical — no retraining!\n'
             'c.add_canonical("XYZlab")',
             font_size=13)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Pipeline Position
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_header(slide, "Where It Fits: ASR Pipeline Integration",
           "Layer 4b — learned fallback after exact lexicon lookup", 4, TOTAL_SLIDES)

add_code_box(slide, Inches(0.6), Inches(1.7), Inches(6.0), Inches(5.0),
             'Audio Input (Urdu speech)\n'
             '       │\n'
             '       ▼\n'
             'Qwen3-ASR ──► Hindi Devanagari\n'
             '       │\n'
             '       ▼\n'
             'Layer 1: Phoneme Mapping (char-by-char)\n'
             '       │\n'
             '       ▼\n'
             'Layer 2: Vowel Normalisation (regex)\n'
             '       │\n'
             '       ▼\n'
             'Layer 3: Exact Lexicon Lookup (WORD_MAP)\n'
             '       │\n'
             '  found? ──YES──► return canonical\n'
             '       │ NO\n'
             '       ▼\n'
             'Layer 4b: PHONETIC CONTRASTIVE MODEL  ◄──\n'
             '       │                                  │\n'
             '       ▼                            PHONETIC=1\n'
             'Final Roman Urdu Output',
             font_size=13)

add_bullet_frame(slide, Inches(7.2), Inches(1.7), Inches(5.5), Inches(4.5), [
    "▸ Activated by env var: PHONETIC=1",
    "▸ NEVER overrides an exact lexicon match",
    "▸ Only processes words the exact lexicon missed",
    "▸ Priority chain:",
    "   1. Exact lexicon (highest priority)",
    "   2. Phonetic model (learned fallback)",
    "   3. Resolver (deprecated, rule-based)",
    "   4. Unchanged (no correction available)",
    "",
    "▸ Threshold configurable via PHONETIC_THRESHOLD",
    "  (default: 0.90)",
], font_size=15, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Model Architecture Overview
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_header(slide, "Model Architecture: CharEncoder",
           "Character-level bidirectional GRU with dual pooling", 5, TOTAL_SLIDES)

add_code_box(slide, Inches(0.6), Inches(1.7), Inches(6.0), Inches(5.0),
             'Input: "chugataai"  (string)\n'
             '  ↓ Vocab.encode()  — lowercase, char→id\n'
             '[3, 8, 21, 7, 1, 20, 1, 1, 9]  (char ids)\n'
             '  ↓\n'
             'nn.Embedding(~30 × 96)    → (B, T, 96)\n'
             '  ↓\n'
             'Bidirectional GRU         → (B, T, 256)\n'
             '  2 layers, hidden=128      128×2 dirs\n'
             '  ↓\n'
             '┌─────────┴──────────┐\n'
             'Masked MEAN    Masked MAX\n'
             '(B, 256)       (B, 256)\n'
             '└─────────┬──────────┘\n'
             '  ↓ concatenate\n'
             '(B, 512)\n'
             '  ↓\n'
             'Linear(512→128) + LayerNorm(128)\n'
             '  ↓\n'
             'L2 Normalise   → (B, 128) unit-norm',
             font_size=13)

# Hyperparameters table
add_text_box(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(0.4),
             "Hyperparameters", font_size=20, bold=True, color=ACCENT_BLUE)

params = [
    ("Parameter", "Value", "Description"),
    ("emb_dim", "96", "Character embedding dimension"),
    ("hidden", "128", "GRU hidden size (per direction)"),
    ("out_dim", "128", "Final embedding dimension"),
    ("num_layers", "2", "Stacked GRU layers"),
    ("dropout", "0.2", "Between GRU layers"),
    ("pooled_dim", "512", "128 × 2 (bi) × 2 (mean+max)"),
    ("vocab_size", "~30", "a-z + <pad> + <unk>"),
]
y = Inches(2.1)
for name, val, desc in params:
    is_header = name == "Parameter"
    c = ACCENT_BLUE if is_header else WHITE
    add_text_box(slide, Inches(7.2), y, Inches(1.6), Inches(0.32),
                 name, font_size=12, bold=is_header, color=c, font_name="Menlo" if not is_header else "Calibri")
    add_text_box(slide, Inches(8.8), y, Inches(0.8), Inches(0.32),
                 val, font_size=12, bold=is_header, color=ACCENT_GREEN if not is_header else c)
    add_text_box(slide, Inches(9.6), y, Inches(3.2), Inches(0.32),
                 desc, font_size=11, color=LIGHT_GRAY if not is_header else c)
    y += Inches(0.35)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Architecture Deep Dive: Why Each Component
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_header(slide, "Architecture Deep Dive: Design Decisions",
           "Every component choice is deliberate", 6, TOTAL_SLIDES)

decisions = [
    ("Character-level (not word-level)",
     "The signal is phonetic — individual character substitutions (z↔j, aa↔a). "
     "Word embeddings would miss sub-word phonetic patterns."),
    ("Bidirectional GRU (not Transformer)",
     "Inputs are short (~8 chars). GRU's sequential bias is natural for characters. "
     "Transformers' O(T²) attention is unnecessary overhead."),
    ("Mean + Max dual pooling",
     "Mean captures average character pattern (robust to noise). "
     "Max captures most salient character features (discriminative n-grams)."),
    ("LayerNorm on projection",
     "Stabilises embedding scale across diverse inputs. Without it, "
     "some inputs produce larger embeddings, biasing similarity."),
    ("L2 normalisation",
     "Places all embeddings on a unit hypersphere. "
     "Cosine similarity = dot product — makes nearest-neighbor a simple matmul."),
    ("Siamese (shared weights)",
     "Same encoder for variants AND canonicals. Ensures they live in "
     "the same space. Half the parameters of a dual-encoder."),
]

y = Inches(1.7)
for title, desc in decisions:
    add_text_box(slide, Inches(0.6), y, Inches(4.0), Inches(0.35),
                 f"▸ {title}", font_size=15, bold=True, color=ACCENT_GREEN)
    add_text_box(slide, Inches(4.8), y, Inches(8.0), Inches(0.35),
                 desc, font_size=13, color=LIGHT_GRAY)
    y += Inches(0.55) if decisions.index((title, desc)) < 4 else Inches(0.5)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Training Data
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_header(slide, "Training Data: Lexicon-Based Pairs",
           "data/lexicons_v2.json — curated canonical → variant mappings", 7, TOTAL_SLIDES)

add_code_box(slide, Inches(0.6), Inches(1.7), Inches(5.5), Inches(2.5),
             '// data/lexicons_v2.json\n'
             '{\n'
             '  "lexicons": {\n'
             '    "lexicon": {\n'
             '      "Siddiqui": ["siddiqee","siddiqi",\n'
             '                   "siddique","sidiqui"],\n'
             '      "Chughtai": ["chugataai","chugatai",\n'
             '                   "chughtaai"],\n'
             '      "operation": ["apareshan","opreshan"],\n'
             '      ...\n'
             '    }\n'
             '  }\n'
             '}',
             font_size=13)

add_text_box(slide, Inches(6.8), Inches(1.5), Inches(5.5), Inches(0.4),
             "Data Splitting (seed=13, reproducible)", font_size=18, bold=True, color=ACCENT_BLUE)

splits = [
    ("Split", "Fraction", "Purpose"),
    ("Train", "~72%", "Pairs the model trains on"),
    ("Validation", "~8%", "Early stopping (carved from train)"),
    ("Held-out", "~20%", "Generalisation test — NEVER seen"),
]
y = Inches(2.1)
for name, frac, purpose in splits:
    is_header = name == "Split"
    c = ACCENT_BLUE if is_header else WHITE
    add_text_box(slide, Inches(6.8), y, Inches(1.5), Inches(0.35),
                 name, font_size=14, bold=is_header, color=c)
    add_text_box(slide, Inches(8.3), y, Inches(1.0), Inches(0.35),
                 frac, font_size=14, color=ACCENT_GREEN if not is_header else c)
    add_text_box(slide, Inches(9.3), y, Inches(3.5), Inches(0.35),
                 purpose, font_size=13, color=LIGHT_GRAY if not is_header else c)
    y += Inches(0.4)

add_bullet_frame(slide, Inches(6.8), Inches(4.0), Inches(5.8), Inches(2.5), [
    "▸ Only single-word entries used (no phrases)",
    "▸ Every canonical keeps ≥1 variant in train",
    "   → canonical stays 'known' while we test unseen spellings",
    "▸ Vocab built from ALL strings (train + heldout)",
    "   → no unknown characters at test time",
], font_size=14, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — InfoNCE Training
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_header(slide, "Training: InfoNCE Contrastive Loss",
           "Pull variants toward their canonical, push away from all others", 8, TOTAL_SLIDES)

add_code_box(slide, Inches(0.6), Inches(1.7), Inches(6.0), Inches(3.5),
             'For each mini-batch of (variant, canonical) pairs:\n\n'
             '1. Encode variants (anchors):\n'
             '   v_emb = encoder(variants)      → (B, 128)\n\n'
             '2. Build candidate set:\n'
             '   candidates = unique_positives + 256 random negatives\n'
             '   c_emb = encoder(candidates)    → (U+256, 128)\n'
             '   ALL encoded FRESH (no stale bank)\n\n'
             '3. Compute InfoNCE loss:\n'
             '   logits = v_emb @ c_emb.T / τ   (τ = 0.07)\n'
             '   loss = CrossEntropy(logits, target)\n\n'
             '4. Optimise:\n'
             '   AdamW (lr=1e-3) + grad clip (5.0)',
             font_size=13)

add_text_box(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(0.4),
             "Key Design Choices", font_size=20, bold=True, color=ACCENT_BLUE)

add_bullet_frame(slide, Inches(7.2), Inches(2.1), Inches(5.5), Inches(4.5), [
    "▸ Unique canonicals as candidates:",
    "  Avoids false negatives when two variants",
    "  share a canonical (e.g. siddiqee + siddiqi → Siddiqui)",
    "",
    "▸ Fresh encoding of negatives every batch:",
    "  All embeddings use current weights — no stale-bank",
    "  collapse that plagues momentum banks",
    "",
    "▸ Temperature τ = 0.07:",
    "  Sharpens the softmax — makes the model decisive",
    "",
    "▸ CosineAnnealingLR scheduler:",
    "  Gradual warmdown prevents late-training instability",
    "",
    "▸ Early stopping (patience=5):",
    "  Monitored on val_recall (top-1 accuracy)",
], font_size=14, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Checkpoint & Index
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_header(slide, "Checkpoint: Self-Contained Single File",
           "Everything to reconstruct the corrector — no external state", 9, TOTAL_SLIDES)

add_code_box(slide, Inches(0.6), Inches(1.7), Inches(6.0), Inches(4.5),
             'phonetic_contrastive_v1.pt  (~3.2 MB)\n'
             '├── state_dict          Model weights (CPU)\n'
             '├── itos                Character vocabulary list\n'
             '├── config\n'
             '│   ├── emb: 96\n'
             '│   ├── hidden: 128\n'
             '│   ├── out: 128\n'
             '│   ├── layers: 2\n'
             '│   └── temp: 0.07\n'
             '├── canonicals          All canonical strings\n'
             '├── canonical_embeddings  (N, 128) pre-computed\n'
             '└── meta\n'
             '    ├── max_epochs\n'
             '    ├── patience\n'
             '    ├── seed\n'
             '    ├── train_pairs\n'
             '    └── best_val_recall',
             font_size=14)

add_text_box(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(0.4),
             "Why Pre-Compute the Index?", font_size=20, bold=True, color=ACCENT_BLUE)

add_bullet_frame(slide, Inches(7.2), Inches(2.1), Inches(5.5), Inches(4.5), [
    "▸ At training end, ALL canonicals are embedded once:",
    "  canon_emb = embed_all(model, canonicals)  → (N, 128)",
    "",
    "▸ Stored in the checkpoint alongside weights",
    "",
    "▸ At inference, loaded directly — no re-encoding",
    "  of the gazetteer at startup",
    "",
    "▸ Inference cost per word:",
    "  1 forward pass (encode query)",
    "  + 1 matmul (query × index.T)",
    "  = O(T + N)  where T = word length, N = canonicals",
    "",
    "▸ Total model load time: < 0.5 seconds",
    "▸ Per-word latency: < 1 ms on CPU",
], font_size=14, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Inference Pipeline
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_header(slide, "Inference: resolve_word() — The Decision Path",
           "Guard → Encode → Compare → Threshold Gate → Return", 10, TOTAL_SLIDES)

add_code_box(slide, Inches(0.6), Inches(1.7), Inches(6.0), Inches(5.0),
             'resolve_word("chugataai"):\n\n'
             '  1. Guard: is_alpha?  ✓\n'
             '  2. Guard: already canonical?  ✗ (not in known set)\n'
             '  3. Guard: len ≥ 3?  ✓ (9 chars)\n\n'
             '  4. Encode query:\n'
             '     q = model("chugataai")  → (1, 128) unit-norm\n\n'
             '  5. Cosine similarity:\n'
             '     sims = q @ index.T      → (N,) scores\n\n'
             '  6. Best match:\n'
             '     best_idx = argmax(sims)  → idx of "Chughtai"\n'
             '     best_score = 0.95\n\n'
             '  7. Threshold gate:\n'
             '     0.95 ≥ 0.90?  ✓\n\n'
             '  8. Return: "Chughtai"  ✓',
             font_size=13)

add_text_box(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(0.4),
             "Decision Outcomes", font_size=20, bold=True, color=ACCENT_BLUE)

outcomes = [
    ("Input", "Score", "Action", "Reason"),
    ("chugataai", "0.95", "→ Chughtai", "Confident match"),
    ("siddiqee", "0.96", "→ Siddiqui", "Confident match"),
    ("apareshan", "0.78", "→ apareshan", "ABSTAIN (< 0.90)"),
    ("area", "—", "→ area", "Already canonical"),
    ("ka", "—", "→ ka", "Too short (< 3)"),
    ("123", "—", "→ 123", "Non-alphabetic"),
]
y = Inches(2.1)
for inp, score, action, reason in outcomes:
    is_header = inp == "Input"
    c = ACCENT_BLUE if is_header else WHITE
    add_text_box(slide, Inches(7.2), y, Inches(1.5), Inches(0.35),
                 inp, font_size=13, bold=is_header, color=c, font_name="Menlo" if not is_header else "Calibri")
    add_text_box(slide, Inches(8.7), y, Inches(0.7), Inches(0.35),
                 score, font_size=13, color=ACCENT_GREEN if not is_header else c)
    add_text_box(slide, Inches(9.4), y, Inches(1.6), Inches(0.35),
                 action, font_size=13, bold=not is_header, color=ACCENT_GREEN if "→" in action and "apareshan" not in action and "area" not in action and "ka" not in action and "123" not in action else (ACCENT_ORANGE if is_header else LIGHT_GRAY))
    add_text_box(slide, Inches(11.0), y, Inches(2.0), Inches(0.35),
                 reason, font_size=12, color=LIGHT_GRAY if not is_header else c)
    y += Inches(0.4)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Safety Mechanisms
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_header(slide, "Safety: Abstain-Over-Corrupt Philosophy",
           "The model's #1 priority is to never make things worse", 11, TOTAL_SLIDES)

safety_items = [
    ("ABSTAIN THRESHOLD (0.90)",
     "If the nearest canonical's cosine similarity is below 0.90, the word is returned "
     "UNCHANGED. This prevents the model from guessing when it isn't sure. The threshold "
     "was determined by sweep evaluation — 0.90 maximises recall while keeping corruption near zero.",
     ACCENT_RED),
    ("KNOWN-CANONICAL SHORT CIRCUIT",
     "If a word is already in the canonical set, the model skips encoding entirely. This prevents "
     "a canonical from being 'corrected' to a similar-sounding different canonical.",
     ACCENT_ORANGE),
    ("MINIMUM LENGTH GUARD (≥ 3)",
     "Words shorter than 3 characters are never processed. Short words like 'ka', 'is', 'ke' have "
     "too little character signal for reliable matching.",
     ACCENT_BLUE),
    ("STATISTICS TRACKING",
     "The corrector logs counts of: matched, abstain_short, abstain_low, already_canonical. "
     "This enables post-hoc analysis of model behaviour on real traffic.",
     ACCENT_GREEN),
]

y = Inches(1.7)
for title, desc, color in safety_items:
    add_text_box(slide, Inches(0.6), y, Inches(12), Inches(0.35),
                 f"▸ {title}", font_size=17, bold=True, color=color)
    add_text_box(slide, Inches(0.9), y + Inches(0.4), Inches(11.5), Inches(0.6),
                 desc, font_size=14, color=LIGHT_GRAY)
    y += Inches(1.1)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Embedding Space Visualisation
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_header(slide, "Embedding Space: How Words Cluster",
           "128-dim unit hypersphere — phonetically similar words are neighbours", 12, TOTAL_SLIDES)

add_code_box(slide, Inches(0.6), Inches(1.7), Inches(5.5), Inches(4.5),
             '    128-dim unit hypersphere\n\n'
             '         Siddiqui ●\n'
             '        /          \\\n'
             ' siddiqee ●        ● siddiqi\n'
             '              \\    /\n'
             '               ● siddique\n'
             '    (all close — cosine > 0.95)\n\n\n'
             '         Chughtai ●\n'
             '        /          \\\n'
             ' chugataai ●      ● chughtaai\n'
             '              \\\n'
             '               ● chugatai\n'
             '    (all close — cosine > 0.92)\n\n\n'
             '                   ● area\n'
             '    (far from both — cosine < 0.50)\n'
             '    → ABSTAIN, left unchanged',
             font_size=13)

add_text_box(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.4),
             "What the Model Learned", font_size=20, bold=True, color=ACCENT_BLUE)

add_bullet_frame(slide, Inches(7.0), Inches(2.1), Inches(5.5), Inches(4.5), [
    "▸ Phonetic equivalences:",
    "  q/k → same cluster    (Siddiqui/Siddiki)",
    "  aa/a → same cluster   (chugataai/chugatai)",
    "  ee/i → same cluster   (siddiqee/siddiqi)",
    "",
    "▸ The model learns these patterns from data,",
    "  not from hand-crafted substitution rules",
    "",
    "▸ Generalises to UNSEEN combinations:",
    "  If it knows q↔k and ee↔i separately,",
    "  it handles 'siddikee' even if never seen",
    "",
    "▸ Distinct words stay far apart:",
    "  'area' vs 'Arif' — different clusters",
    "  High threshold prevents cross-cluster matches",
], font_size=14, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Extensibility
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_header(slide, "Extensibility: Add Canonicals Without Retraining",
           "The key maintainability advantage over rule-based systems", 13, TOTAL_SLIDES)

add_code_box(slide, Inches(0.6), Inches(1.7), Inches(5.8), Inches(2.0),
             '# Runtime: single new canonical\n'
             'corrector.add_canonical("XYZlab")\n'
             '# → encodes with existing model\n'
             '# → appends to index\n'
             '# → immediately resolves garbles\n\n'
             '# Bulk: extend_canonicals.py\n'
             'python -m phonetic_contrastive_model.extend_canonicals \\\n'
             '    --terms new_entities.txt --save-index',
             font_size=13)

add_text_box(slide, Inches(0.6), Inches(4.2), Inches(5.8), Inches(0.4),
             "extend_canonicals.py Pipeline:", font_size=16, bold=True, color=ACCENT_BLUE)

add_bullet_frame(slide, Inches(0.6), Inches(4.7), Inches(5.8), Inches(2.0), [
    "1. Read new entity terms from text file",
    "2. Add each to the model's index",
    "3. Validate against 80-call benchmark",
    "4. Drop ambiguous terms (wrong captures)",
    "5. Save extended index + updated lexicon",
], font_size=14, color=WHITE)

add_text_box(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(0.4),
             "Why This Works", font_size=20, bold=True, color=ACCENT_BLUE)

add_bullet_frame(slide, Inches(7.2), Inches(2.1), Inches(5.5), Inches(4.5), [
    "▸ The model learned a GENERAL character-level",
    "  embedding space during training",
    "",
    "▸ A new canonical placed in this space will",
    "  'attract' its phonetic variants — even ones",
    "  never seen during training",
    "",
    "▸ rebuild_index() is IDEMPOTENT:",
    "  Re-running always produces the same result",
    "  regardless of prior edits",
    "",
    "▸ Comparison with rule-based approach:",
    "  Rule-based: add N variants × M rules per term",
    "  Contrastive: add_canonical(term) — done",
], font_size=14, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Evaluation & Results
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_header(slide, "Evaluation: The 4-Number Scorecard",
           "Four metrics that comprehensively assess model quality", 14, TOTAL_SLIDES)

metrics = [
    ("1", "HELD-OUT VARIANT RECALL",
     "Unseen spelling variants → correct canonical?\nThe GENERALISATION number.\n"
     "Top-1 (with abstain), Top-1 (raw), Top-3",
     "~97%", ACCENT_GREEN),
    ("2", "ABSTAIN SAFETY",
     "Real gold words (not canonicals) → % left unchanged\nThe ANTI-CORRUPTION number.\n"
     "Tested on benchmark words",
     "~99%", ACCENT_GREEN),
    ("3", "EXACT-NAME HELD-OUT",
     "Entity names (Siddiqui, Chughtai) → exact match\nThe DECISIVE test.\n"
     "Names are highest-value corrections",
     "High", ACCENT_BLUE),
    ("4", "80-CALL diff_words",
     "End-to-end: v2 baseline vs v2 + model\nThe PRODUCTION metric.\n"
     "Accuracy-neutral (fuzzy already forgives drift)",
     "Stable", ACCENT_BLUE),
]

y = Inches(1.7)
for num, title, desc, result, color in metrics:
    add_text_box(slide, Inches(0.6), y, Inches(0.5), Inches(0.35),
                 num, font_size=24, bold=True, color=color)
    add_text_box(slide, Inches(1.2), y, Inches(4.5), Inches(0.35),
                 title, font_size=16, bold=True, color=color)
    lines = desc.split("\n")
    for j, line in enumerate(lines):
        add_text_box(slide, Inches(1.2), y + Inches(0.35 + j * 0.25), Inches(7.0), Inches(0.25),
                     line, font_size=12, color=LIGHT_GRAY)
    add_text_box(slide, Inches(10.5), y + Inches(0.15), Inches(2.0), Inches(0.5),
                 result, font_size=22, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    y += Inches(1.35)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Comparison & Summary
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_header(slide, "Summary: Contrastive Model vs Rule-Based Resolver",
           "", 15, TOTAL_SLIDES)

comparison = [
    ("Aspect", "Rule-Based Resolver", "Phonetic Contrastive Model"),
    ("Held-out recall", "~30%", "~97%"),
    ("Maintenance", "O(N²) — manual rules", "add_canonical() — done"),
    ("Generalisation", "None — only programmed", "Unseen spellings handled"),
    ("Abstain safety", "Manual tuning", "Built-in 0.90 threshold"),
    ("Latency", "~0.01 ms/word", "~0.1 ms/word"),
    ("New canonicals", "Add variants manually", "No retrain needed"),
    ("Model size", "N/A (rules)", "~3.2 MB checkpoint"),
    ("Interpretability", "Rules readable", "Embedding space"),
]

y = Inches(1.6)
for aspect, rule, model_val in comparison:
    is_header = aspect == "Aspect"
    c = ACCENT_BLUE if is_header else WHITE
    add_text_box(slide, Inches(0.6), y, Inches(3.5), Inches(0.38),
                 aspect, font_size=14, bold=is_header, color=c)
    add_text_box(slide, Inches(4.2), y, Inches(4.0), Inches(0.38),
                 rule, font_size=13, bold=is_header,
                 color=c if is_header else ACCENT_RED)
    add_text_box(slide, Inches(8.5), y, Inches(4.5), Inches(0.38),
                 model_val, font_size=13, bold=is_header,
                 color=c if is_header else ACCENT_GREEN)
    y += Inches(0.42)

add_text_box(slide, Inches(0.6), Inches(5.8), Inches(12), Inches(0.8),
             "The Phonetic Contrastive Model delivers 3× the recall, near-zero corruption risk, "
             "and zero-maintenance canonical onboarding — in a 3.2 MB self-contained checkpoint.",
             font_size=18, bold=True, color=ACCENT_ORANGE)


# ═══════════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════════
prs.save(str(OUTPUT))
print(f"Saved: {OUTPUT}")
