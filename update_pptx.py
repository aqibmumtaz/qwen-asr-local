#!/usr/bin/env python3
"""
Update AI-Helpdesk-Presentation.pptx with detailed STT & TTS challenges
from ard/asr-tts-challenges.md.

Replaces slide 12 (Current Challenges) with multiple detailed challenge slides,
preserving slides 1-11 and 13-15 (renumbered).
"""

import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

INPUT = Path("/Users/AqibMumtaz/Downloads/AI-Helpdesk-Presentation.pptx")
OUTPUT = Path("/Users/AqibMumtaz/Downloads/AI-Helpdesk-Presentation-Updated.pptx")

prs = Presentation(str(INPUT))

# Get slide dimensions
slide_width = prs.slide_width
slide_height = prs.slide_height

# Use the blank layout (same as existing slides)
blank_layout = prs.slide_layouts[6]  # Usually "Blank"

# Colors matching existing theme
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)  # Dark navy
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BLUE = RGBColor(0x4E, 0xC5, 0xF1)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_GREEN = RGBColor(0x4E, 0xC9, 0xA0)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)


def add_text_box(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=14,
    bold=False,
    color=WHITE,
    alignment=PP_ALIGN.LEFT,
):
    """Add a text box with specified properties."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf


def add_title_subtitle(slide, title, subtitle, slide_num):
    """Add title + subtitle + slide number in consistent format."""
    # Title
    add_text_box(
        slide,
        Inches(0.5),
        Inches(0.3),
        Inches(9),
        Inches(0.6),
        title,
        font_size=28,
        bold=True,
        color=WHITE,
    )
    # Subtitle
    add_text_box(
        slide,
        Inches(0.5),
        Inches(0.85),
        Inches(9),
        Inches(0.4),
        subtitle,
        font_size=12,
        color=LIGHT_GRAY,
    )
    # Slide number
    add_text_box(
        slide,
        Inches(9.2),
        Inches(0.3),
        Inches(0.5),
        Inches(0.4),
        str(slide_num),
        font_size=11,
        color=LIGHT_GRAY,
        alignment=PP_ALIGN.RIGHT,
    )


def add_bullet_block(slide, left, top, width, height, items, font_size=11, color=WHITE):
    """Add a bulleted text block."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(4)
    return tf


# --- Insert new challenge slides AFTER slide 12 (keep existing slide 12) ---
# New slides will be added at end then reordered to position after slide 12
new_slides = []

# === SLIDE: STT Challenges Overview ===
slide = prs.slides.add_slide(blank_layout)
add_title_subtitle(
    slide,
    "STT & TTS Challenges",
    "CLL Agent — Technical challenges in Speech-to-Text and Text-to-Speech for Urdu telephony",
    "10",
)
# Two-column layout
add_text_box(
    slide,
    Inches(0.5),
    Inches(1.4),
    Inches(4.5),
    Inches(0.4),
    "⚠️  Speech-to-Text (STT)",
    font_size=16,
    bold=True,
    color=ACCENT_ORANGE,
)
add_bullet_block(
    slide,
    Inches(0.5),
    Inches(1.9),
    Inches(4.5),
    Inches(4.5),
    [
        "8 kHz telephony audio vs 16 kHz model requirement",
        "Dual-speaker single channel (agent + caller)",
        "No 'language Urdu' option — Hindi hint drops nuktas",
        "Code-switching: Urdu + English medical terms",
        "Lack of Urdu-specific training data",
        "Regional accent variation (Lahore, Karachi, Peshawar)",
    ],
    font_size=11,
)

add_text_box(
    slide,
    Inches(5.2),
    Inches(1.4),
    Inches(4.5),
    Inches(0.4),
    "🔊  Text-to-Speech (TTS)",
    font_size=16,
    bold=True,
    color=ACCENT_BLUE,
)
add_bullet_block(
    slide,
    Inches(5.2),
    Inches(1.9),
    Inches(4.5),
    Inches(4.5),
    [
        "16 kHz TTS → 8 kHz telephony degrades quality",
        "No high-quality on-prem Urdu TTS models",
        "Missing diacritics (aerab) — ambiguous pronunciation",
        "Mixed Urdu/English pronunciation in responses",
        "Number/date/price normalisation for Urdu",
        "Naturalness & real-time latency requirements",
        "No custom brand voice for Chughtai Lab",
    ],
    font_size=11,
)
new_slides.append(slide)

# === SLIDE: STT Challenge 1 — Audio Quality ===
slide = prs.slides.add_slide(blank_layout)
add_title_subtitle(
    slide,
    "STT: 8 kHz Telephony vs 16 kHz ASR",
    "Audio quality fundamentally limits recognition accuracy",
    "10a",
)

add_text_box(
    slide,
    Inches(0.5),
    Inches(1.4),
    Inches(9),
    Inches(0.8),
    "Problem: Upsampling 8→16 kHz adds no information. High-frequency band (4–8 kHz) "
    "containing Urdu fricatives (ش، س، ز، ف، خ، غ) is permanently lost at codec stage.",
    font_size=12,
    color=LIGHT_GRAY,
)

add_text_box(
    slide,
    Inches(0.5),
    Inches(2.4),
    Inches(4),
    Inches(0.3),
    "Impact:",
    font_size=12,
    bold=True,
    color=ACCENT_ORANGE,
)
add_bullet_block(
    slide,
    Inches(0.5),
    Inches(2.8),
    Inches(4.3),
    Inches(2),
    [
        "Higher WER on fricative-heavy Urdu vocabulary",
        "Muffled signal degrades confidence across the board",
        "Medical terms with sibilants especially affected",
    ],
    font_size=11,
)

add_text_box(
    slide,
    Inches(5.2),
    Inches(2.4),
    Inches(4),
    Inches(0.3),
    "Mitigation:",
    font_size=12,
    bold=True,
    color=ACCENT_GREEN,
)
add_bullet_block(
    slide,
    Inches(5.2),
    Inches(2.8),
    Inches(4.3),
    Inches(2),
    [
        "Request 16 kHz recording from PBX before codec",
        "Use telephony-trained acoustic model variant",
        "Apply spectral enhancement pre-processing",
    ],
    font_size=11,
)
new_slides.append(slide)

# === SLIDE: STT Challenge 2 — Dual Speaker ===
slide = prs.slides.add_slide(blank_layout)
add_title_subtitle(
    slide,
    "STT: Dual-Speaker Single Channel",
    "Agent + Caller mixed on one audio channel",
    "10b",
)

add_text_box(
    slide,
    Inches(0.5),
    Inches(1.4),
    Inches(4.3),
    Inches(0.3),
    "Caller (Patient):",
    font_size=12,
    bold=True,
    color=ACCENT_BLUE,
)
add_bullet_block(
    slide,
    Inches(0.5),
    Inches(1.8),
    Inches(4.3),
    Inches(1.2),
    [
        "Normal volume, variable pace",
        "May have regional accent",
        "Handled reasonably by model",
    ],
    font_size=11,
)

add_text_box(
    slide,
    Inches(5.2),
    Inches(1.4),
    Inches(4.3),
    Inches(0.3),
    "Agent (Helpdesk Staff):",
    font_size=12,
    bold=True,
    color=ACCENT_ORANGE,
)
add_bullet_block(
    slide,
    Inches(5.2),
    Inches(1.8),
    Inches(4.3),
    Inches(1.2),
    [
        "Low volume (distant mic)",
        "Very fast speech (professional cadence)",
        "Domain jargon — under-represented in training",
    ],
    font_size=11,
)

add_text_box(
    slide,
    Inches(0.5),
    Inches(3.2),
    Inches(9),
    Inches(0.3),
    "Impact:",
    font_size=12,
    bold=True,
    color=ACCENT_ORANGE,
)
add_bullet_block(
    slide,
    Inches(0.5),
    Inches(3.6),
    Inches(9),
    Inches(1.5),
    [
        "Agent utterances get higher error rate; fast speech causes word merging",
        "No speaker separation — ASR sees both as one signal",
        "Low-volume agent audio partially drowned by caller",
    ],
    font_size=11,
)

add_text_box(
    slide,
    Inches(0.5),
    Inches(4.8),
    Inches(9),
    Inches(0.3),
    "Mitigation: Separate recording channels; VAD + volume normalisation; diarisation before ASR",
    font_size=11,
    color=ACCENT_GREEN,
)
new_slides.append(slide)

# === SLIDE: STT Challenge 3 — Language Hint Dilemma ===
slide = prs.slides.add_slide(blank_layout)
add_title_subtitle(
    slide,
    "STT: Language Hint Dilemma",
    "No 'language Urdu' option — neither Hindi nor English hint is correct",
    "10c",
)

add_text_box(
    slide,
    Inches(0.5),
    Inches(1.5),
    Inches(4.3),
    Inches(0.3),
    "language=Hindi",
    font_size=14,
    bold=True,
    color=ACCENT_BLUE,
)
add_bullet_block(
    slide,
    Inches(0.5),
    Inches(1.9),
    Inches(4.3),
    Inches(1.8),
    [
        "Outputs Devanagari — preserves Urdu words",
        "DROPS nuktas (़): क़→क, ज़→ज, फ़→फ",
        "قابل→kaabil not qaabil",
        "زندگی→jindagi not zindagi",
    ],
    font_size=11,
)

add_text_box(
    slide,
    Inches(5.2),
    Inches(1.5),
    Inches(4.3),
    Inches(0.3),
    "language=English",
    font_size=14,
    bold=True,
    color=ACCENT_ORANGE,
)
add_bullet_block(
    slide,
    Inches(5.2),
    Inches(1.9),
    Inches(4.3),
    Inches(1.8),
    [
        "Preserves nuktas — more accurate phonetically",
        "But transcribes Urdu phrases in Latin/English",
        "Entire segments can come out wrong script",
        "Inconsistent behaviour on ambiguous audio",
    ],
    font_size=11,
)

add_text_box(
    slide,
    Inches(0.5),
    Inches(3.9),
    Inches(9),
    Inches(0.8),
    "Root Cause: Urdu is not in model's supported language list. Urdu uses Perso-Arabic script; "
    "model transcribes as Hindi/Devanagari. Nuktas distinguish Perso-Arabic phonemes but Hindi training data omits them.",
    font_size=11,
    color=LIGHT_GRAY,
)

add_text_box(
    slide,
    Inches(0.5),
    Inches(4.8),
    Inches(9),
    Inches(0.5),
    "Mitigation: Fine-tune on nukta-annotated Devanagari; post-process with rule-based nukta restoration; "
    "custom lexicon for domain vocabulary",
    font_size=11,
    color=ACCENT_GREEN,
)
new_slides.append(slide)

# === SLIDE: STT Challenge 4+5 — Training Challenges ===
slide = prs.slides.add_slide(blank_layout)
add_title_subtitle(
    slide,
    "STT: Model Training Challenges",
    "Why off-the-shelf ASR fails for Urdu telephony",
    "10d",
)

add_text_box(
    slide,
    Inches(0.5),
    Inches(1.4),
    Inches(4.3),
    Inches(0.3),
    "Data Scarcity:",
    font_size=12,
    bold=True,
    color=ACCENT_ORANGE,
)
add_bullet_block(
    slide,
    Inches(0.5),
    Inches(1.8),
    Inches(4.3),
    Inches(2.2),
    [
        "Urdu severely under-represented in Whisper/Qwen3",
        "No large transcribed Urdu telephony corpus exists",
        "CommonVoice Urdu: ~100hrs read-speech only",
        "Domain annotation requires native expertise",
        "Medical vocabulary completely out-of-distribution",
    ],
    font_size=10,
)

add_text_box(
    slide,
    Inches(5.2),
    Inches(1.4),
    Inches(4.3),
    Inches(0.3),
    "Fine-tuning Complexity:",
    font_size=12,
    bold=True,
    color=ACCENT_BLUE,
)
add_bullet_block(
    slide,
    Inches(5.2),
    Inches(1.8),
    Inches(4.3),
    Inches(2.2),
    [
        "Requires GPU infrastructure + careful scheduling",
        "Risk of catastrophic forgetting on other languages",
        "Ground truth must use consistent Nastaliq orthography",
        "Evaluation needs native Urdu annotators",
        "Regional accents (Lahore/Karachi/Peshawar) vary",
    ],
    font_size=10,
)

add_text_box(
    slide,
    Inches(0.5),
    Inches(4.2),
    Inches(9),
    Inches(0.3),
    "Domain Mismatch:",
    font_size=12,
    bold=True,
    color=ACCENT_ORANGE,
)
add_bullet_block(
    slide,
    Inches(0.5),
    Inches(4.5),
    Inches(9),
    Inches(1.5),
    [
        "Models trained on news/podcasts — not live telephony with interruptions, hesitations, noise",
        "Medical vocabulary (haematology, biochemistry, radiology) is OOV in all general models",
        "Code-switching (Urdu + English medical terms + Punjabi) not in training distribution",
    ],
    font_size=10,
)
new_slides.append(slide)

# === SLIDE: STT Solution — Hindi → Roman Urdu Pipeline ===
slide = prs.slides.add_slide(blank_layout)
add_title_subtitle(
    slide,
    "STT Solution: Hindi → Roman Urdu Pipeline",
    "Custom post-processing layer built at BitLogix bridges the ASR gap",
    "10e",
)

add_text_box(
    slide,
    Inches(0.5),
    Inches(1.4),
    Inches(9),
    Inches(0.3),
    "Pipeline Architecture:",
    font_size=12,
    bold=True,
    color=ACCENT_GREEN,
)
add_bullet_block(
    slide,
    Inches(0.5),
    Inches(1.8),
    Inches(9),
    Inches(2.5),
    [
        "Layer 1: ASR → Devanagari (Qwen3-ASR outputs Hindi for Urdu speech)",
        "Layer 2: Phoneme mapping — Devanagari chars → Roman Urdu (nukta-aware: फ़→f, ज़→z, क़→q)",
        "Layer 3: Schwa deletion — removes Sanskrit inherent vowel silent in Urdu (करम→karam not karama)",
        "Layer 4: Lexicon correction — 985 word corrections + 286 proper nouns (medical, names, branches)",
    ],
    font_size=11,
)

add_text_box(
    slide,
    Inches(0.5),
    Inches(3.8),
    Inches(9),
    Inches(0.3),
    "What It Solves:",
    font_size=12,
    bold=True,
    color=ACCENT_GREEN,
)
add_bullet_block(
    slide,
    Inches(0.5),
    Inches(4.1),
    Inches(9),
    Inches(2),
    [
        "No native Urdu ASR → Devanagari re-mapped to Roman Urdu for downstream NLP",
        "Nukta omission → English hint + nukta-aware phoneme map handles ف़/ज़/ق correctly",
        "Medical/proper nouns → dictionary covers lab branches, doctor names, test names",
        "Script mismatch → clean Roman Urdu output compatible with intent classifiers & LLM",
    ],
    font_size=11,
)

add_text_box(
    slide,
    Inches(0.5),
    Inches(5.5),
    Inches(9),
    Inches(0.5),
    "Limitation: Lexicon must be maintained manually; regional accents can miss the lexicon; "
    "nukta emission remains inconsistent across audio quality levels.",
    font_size=10,
    color=LIGHT_GRAY,
)
new_slides.append(slide)

# === SLIDE: TTS Challenges ===
slide = prs.slides.add_slide(blank_layout)
add_title_subtitle(
    slide,
    "TTS Challenges",
    "Text-to-Speech barriers for production Urdu voice responses",
    "10f",
)

add_text_box(
    slide,
    Inches(0.5),
    Inches(1.4),
    Inches(4.3),
    Inches(0.3),
    "Quality & Availability:",
    font_size=12,
    bold=True,
    color=ACCENT_ORANGE,
)
add_bullet_block(
    slide,
    Inches(0.5),
    Inches(1.8),
    Inches(4.3),
    Inches(2.5),
    [
        "16 kHz TTS → 8 kHz telephony: loses naturalness",
        "Urdu fricatives (ش، س، ز، خ) especially degraded",
        "No high-quality on-prem Urdu TTS exists",
        "Cloud options (Google/Azure) excluded by policy",
        "Fine-tuning needs 5–20 hrs professional recording",
        "Missing aerab (diacritics) = ambiguous pronunciation",
    ],
    font_size=10,
)

add_text_box(
    slide,
    Inches(5.2),
    Inches(1.4),
    Inches(4.3),
    Inches(0.3),
    "Production Requirements:",
    font_size=12,
    bold=True,
    color=ACCENT_BLUE,
)
add_bullet_block(
    slide,
    Inches(5.2),
    Inches(1.8),
    Inches(4.3),
    Inches(2.5),
    [
        "Mixed Urdu/English pronunciation in responses",
        "Number/date normalisation (Rs. 1,450 → Roman Urdu)",
        "Real-time latency: <1s response time needed",
        "Empathetic, clear delivery for patient trust",
        "Custom Chughtai Lab brand voice identity",
        "Consistency across sessions and call types",
    ],
    font_size=10,
)

# Model comparison table as text
add_text_box(
    slide,
    Inches(0.5),
    Inches(4.5),
    Inches(9),
    Inches(0.3),
    "On-Prem Model Options:",
    font_size=12,
    bold=True,
    color=WHITE,
)
add_text_box(
    slide,
    Inches(0.5),
    Inches(4.9),
    Inches(9),
    Inches(1.2),
    "HMM/Concatenative: CPU-only, robotic, fast  |  VITS (fine-tuned): mid GPU, best quality/speed  |  "
    "XTTS v2: 8GB VRAM, high quality but slow  |  Diffusion TTS: not viable for real-time",
    font_size=10,
    color=LIGHT_GRAY,
)
new_slides.append(slide)

# --- Now reorder slides: move new slides (at end) to after slide 12 ---
sldIdLst = prs.slides._sldIdLst
total = len(sldIdLst)
n_new = len(new_slides)

# Current order: [0..14] (original 15 slides) + [15..15+n_new-1] (new slides at end)
# We want: [0..11] (slides 1-12) + [15..15+n_new-1] (new) + [12..14] (original slides 13-15)
n_before = 12  # keep slides 1-12 (indices 0-11)
n_after_old = total - n_new - n_before  # original slides 13-15 (indices 12-14)

all_ids = list(sldIdLst)
for sid in all_ids:
    sldIdLst.remove(sid)

reordered = (
    all_ids[:n_before]
    + all_ids[n_before + n_after_old :]
    + all_ids[n_before : n_before + n_after_old]
)
for sid in reordered:
    sldIdLst.append(sid)

prs.save(str(OUTPUT))
print(f"✓ Saved updated presentation: {OUTPUT}")
print(f"  Total slides: {total}")
print(f"  New challenge slides inserted after slide 12: {n_new}")
print(
    f"  Slides 1-12: preserved | Slides 13-{12+n_new}: STT/TTS challenges | Slides {13+n_new}-{total}: original 13-15"
)
